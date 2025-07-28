from pptx import Presentation

def read_pptx(file_path):
    lines = []
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        lines.append(text)

            elif shape.shape_type == 19:  # 19 = TABLE
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        # Combine row into one line so regex can hit entire row
                        lines.append(" | ".join(row_text))

    return lines
